from flask import Flask, request, render_template, send_from_directory, send_file
import PyPDF2
import os
import asyncio
import edge_tts
from pydub import AudioSegment
from llama_index.llms.groq import Groq
from llama_index.core.llms import ChatMessage
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re
import textwrap

app = Flask(__name__)

# Ensure required directories exist
os.makedirs('static', exist_ok=True)
os.makedirs('uploads', exist_ok=True)
os.makedirs('generated', exist_ok=True)

# Initialize LLM (Groq model)
llm = Groq(model="llama3-70b-8192", api_key="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")

def extract_sections_from_pdf(pdf_path):
    sections = {}
    try:
        with open(pdf_path, "rb") as file:
            reader = PyPDF2.PdfReader(file)
            text = "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])

        if not text:
            return {}

        # Improved regex to capture more sections based on your Python script
        pattern = re.compile(r'(\n(?:[IVXLCDM]\.*\s*)?[A-Z][A-Z\s\-0-9]+)\n')  # Adjust pattern for additional sections
        matches = list(pattern.finditer(text))

        if not matches:
            return {}

        # Improved logic to extract more sections and handle cases where sections might be missed
        for i in range(len(matches)):
            start = matches[i].end()
            end = matches[i+1].start() if i+1 < len(matches) else len(text)
            section_name = matches[i].group(1).strip()
            section_content = text[start:end].strip()
            sections[section_name] = section_content

        return sections
    except Exception as e:
        print(f"Error extracting PDF: {str(e)}")
        return {}


# Function to extract text from PDF
def extract_text_from_pdf(pdf_path):
    with open(pdf_path, "rb") as file:
        reader = PyPDF2.PdfReader(file)
        return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])

# Function to generate informal podcast script
def generate_informal_podcast_script(extracted_text, target_length=2500):
    if len(extracted_text) > target_length:
        extracted_text = extracted_text[:target_length]

    messages = [
        ChatMessage(role="system", content="""You are an AI that converts research paper text into a simple, fun podcast script.
        The podcast should be lively, with a casual, friendly tone. Keep it light, fun, and easy to understand.
        The script should involve one person explaining things in a simple way."""),
        ChatMessage(role="user", content=extracted_text),
    ]

    response = llm.chat(messages)
    return response.message.content

# Async function to generate speech
def generate_speech(text, voice, filename):
    async def generate():
        tts = edge_tts.Communicate(text=text, voice=voice)
        await tts.save(filename)
    asyncio.run(generate())

# Function to generate podcast audio
def generate_podcast_audio(script_file, output_filename="static/podcast.mp3", voice="en-GB-SoniaNeural"):
    with open(script_file, "r", encoding="utf-8") as file:
        script_text = file.read()

    generate_speech(script_text, voice, output_filename)
    return output_filename

@app.route("/", methods=["GET"])
def index():
    return render_template("index.html")

# Function to generate informal summaries using Llama
def generate_informal_summary(text):
    prompt = "Summarize the text into 3 very simple points that anyone can understand."

    messages = [
        ChatMessage(role="system", content=prompt),
        ChatMessage(role="user", content=text)
    ]
    try:
        response = llm.chat(messages).message.content.strip()
        return response
    except Exception as e:
        print(f"Error generating summary: {str(e)}")
        return ""
# Function to extract sections from PDF (updated for efficiency)

# Function to create PPT (with informal presentation style)
def create_informal_ppt(summaries, ppt_type, output_path):
    ppt = Presentation()
    title_slide_layout = ppt.slide_layouts[0]
    content_slide_layout = ppt.slide_layouts[1]

    # Title Slide
    slide = ppt.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"{ppt_type.capitalize()} Summary Presentation"
    subtitle.text = "Automatically Generated from Research Paper"

    def wrap_text(text, max_chars=80):
        return "\n".join(textwrap.wrap(text, width=max_chars))

    for section, summary in summaries.items():
        slide = ppt.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = section
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 69, 0)  # Bright color for informal tone

        # Add simple background color
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 0)  # Light yellow background

        # Add image (optional)
        try:
            img_path = 'static/sample_image.jpg'  # You can use any image path here
            slide.shapes.add_picture(img_path, 100, 100, width=300, height=200)
        except Exception as e:
            print(f"Error adding image: {str(e)}")

        bullet_points = summary.split("\n")
        tf = content.text_frame
        tf.clear()

        for point in bullet_points:
            wrapped_point = wrap_text(point.strip(), 70)
            p = tf.add_paragraph()
            p.text = wrapped_point
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0, 0, 255)  # Blue text for clarity
            p.line_spacing = Pt(14)

    ppt.save(output_path)

# Function to process PDF for informal PPT
def process_pdf_for_informal_ppt(pdf_path, ppt_type):
    sections = extract_sections_from_pdf(pdf_path)
    summaries = {section: generate_informal_summary(content) for section, content in sections.items() if section != "REFERENCES"}
    output_path = os.path.join("generated", f"{ppt_type}_summary_presentation.pptx")
    create_informal_ppt(summaries, ppt_type, output_path)
    return output_path

@app.route("/ppt", methods=["GET", "POST"])
def ppt():
    if request.method == "POST":
        if "file" not in request.files:
            return "No file uploaded", 400
        file = request.files["file"]
        if file.filename == "":
            return "No selected file", 400
        ppt_type = request.form.get("pptType", "informal")
        file_path = os.path.join("uploads", file.filename)
        file.save(file_path)
        output_ppt_path = process_pdf_for_informal_ppt(file_path, ppt_type)
        return render_template("ppt.html", ppt_generated=True, ppt_path=output_ppt_path)
    return render_template("ppt.html", ppt_generated=False)

@app.route("/download/ppt")
def download_ppt():
    ppt_path = os.path.join("generated", "informal_summary_presentation.pptx")  # Defaulting to informal
    if not os.path.exists(ppt_path):
        return "No PPT generated yet", 404
    return send_file(ppt_path, as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True)
